"""
Build a polarity map PowerPoint from the template.

Uses python-pptx to open the template and fill in the
12 placeholder text boxes with data from the LLM's JSON output.
"""
from pathlib import Path
from typing import List
from pptx import Presentation
from pptx.util import Pt
from lxml import etree


# Mapping from our JSON keys to placeholder indices in the template
PLACEHOLDER_MAP = {
    "gps":               14,   # Top center - Übergeordnetes Ziel
    "deeper_fear":       15,   # Bottom center - Schlimmster Fall
    "pole_a":            16,   # Center left - Pol 1
    "pole_b":            17,   # Center right - Pol 2
    "upsides_a":         18,   # Top-left inner quadrant
    "downsides_a":       19,   # Bottom-left inner quadrant
    "upsides_b":         20,   # Top-right inner quadrant
    "downsides_b":       21,   # Bottom-right inner quadrant
    "action_steps_a":    22,   # Far left top
    "early_warnings_a":  23,   # Far left bottom
    "action_steps_b":    24,   # Far right top
    "early_warnings_b":  25,   # Far right bottom
}

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _make_paragraph_xml(text: str, lang: str = "de-DE", font_size: int = 1100) -> etree._Element:
    """Create a single <a:p> element with a run containing the given text."""
    p = etree.SubElement(etree.Element("dummy"), f"{{{NS}}}p")
    pPr = etree.SubElement(p, f"{{{NS}}}pPr")
    pPr.set("lvl", "0")
    r = etree.SubElement(p, f"{{{NS}}}r")
    rPr = etree.SubElement(r, f"{{{NS}}}rPr")
    rPr.set("lang", lang)
    rPr.set("dirty", "0")
    rPr.set("sz", str(font_size))
    t = etree.SubElement(r, f"{{{NS}}}t")
    t.text = text
    return p


def _set_placeholder_text(slide, idx: int, text: str, font_size: int = 1200) -> None:
    """Set a single text value on a placeholder by its idx."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            txBody = tf._txBody

            # Remove all existing paragraphs
            for p in txBody.findall(f"{{{NS}}}p"):
                txBody.remove(p)

            # Add one paragraph with the text
            new_p = _make_paragraph_xml(text, font_size=font_size)
            txBody.append(new_p)
            return


def _set_placeholder_list(slide, idx: int, items: List[str], font_size: int = 1100) -> None:
    """Set a numbered list on a placeholder, one paragraph per item."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            txBody = tf._txBody

            # Remove all existing paragraphs
            for p in txBody.findall(f"{{{NS}}}p"):
                txBody.remove(p)

            # Add one paragraph per item
            for item in items:
                new_p = _make_paragraph_xml(item, font_size=font_size)
                txBody.append(new_p)
            return


def build_powerpoint(polarity_data: dict, template_path: Path, output_path: Path) -> Path:
    """
    Fill the polarity map template with data from the LLM.

    Args:
        polarity_data: Dict with keys matching PLACEHOLDER_MAP
        template_path: Path to the .pptx template
        output_path: Path for the output .pptx

    Returns:
        Path to the generated .pptx file
    """
    prs = Presentation(str(template_path))
    slide = prs.slides[0]

    # Single-value fields (GPS, Deeper Fear, Pole names)
    single_fields = ["gps", "deeper_fear", "pole_a", "pole_b"]
    for field in single_fields:
        if field in polarity_data and polarity_data[field]:
            idx = PLACEHOLDER_MAP[field]
            _set_placeholder_text(slide, idx, polarity_data[field])

    # List fields (quadrants, action steps, early warnings)
    list_fields = [
        "upsides_a", "upsides_b",
        "downsides_a", "downsides_b",
        "action_steps_a", "action_steps_b",
        "early_warnings_a", "early_warnings_b",
    ]
    for field in list_fields:
        if field in polarity_data and polarity_data[field]:
            idx = PLACEHOLDER_MAP[field]
            _set_placeholder_list(slide, idx, polarity_data[field])

    prs.save(str(output_path))
    return output_path